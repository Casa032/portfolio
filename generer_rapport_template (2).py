def _demander_verdict(titre: str, texte: str, bloc_few_shot: str = "", bloc_estimation: str = ""):
    ...
    systeme = (
        "..."
        f"{exemples_txt}"
        f"{bloc_estimation}"
        "..."
    )


"Réponds UNIQUEMENT par un objet JSON valide, sans texte autour :\n"
'{"score": <entier 0-10>, "raison": "1 phrase précise citant le critère '
'déterminant : présence/absence de sanction ou décision officielle, '
'portée contraignante ou non, lien direct ou superficiel avec le '
'périmètre DPO", "resume_structure": "le résumé structuré multi-lignes"}'

def scorer_llm_article(article: dict) -> dict:
    """ Ajoute llm_score, llm_raison et llm_resume à l'article, en place.
        Déclenche un arbitrage automatique si le score du LLM diverge trop
        de l'estimation statistique k-NN. """
    if not LLM_ACTIVE:
        return article

    texte = (
        article.get("contenu")
        or article.get("resume")
        or article.get("llm_resume")
        or ""
    )

    exemples = fs.selectionner_exemples(texte, k=5, min_par_tranche=1)
    bloc_few_shot = fs.formatter_exemples_pour_prompt(exemples)

    verdict = _demander_verdict(article.get("titre"), texte, bloc_few_shot)
    if verdict is None:
        return article

    score, raison, resume_struct = verdict
    article["llm_score"], article["llm_raison"], article["llm_resume"] = score, raison, resume_struct

    estimation = fs.estimer_score_knn(texte, k=10)
    if estimation is not None:
        ecart = abs(score - estimation["score_estime"])
        article["score_estime_knn"] = estimation["score_estime"]
        article["ecart_knn"] = ecart

        if ecart > SEUIL_ECART_ARBITRAGE:
            arbitrage = _demander_arbitrage(
                article.get("titre"), texte, score, raison,
                estimation["score_estime"], bloc_few_shot,
            )
            if arbitrage is not None:
                score_final, justification_arbitrage = arbitrage
                article["llm_score"] = score_final
                article["arbitrage_effectue"] = True
                article["justification_arbitrage"] = justification_arbitrage
            else:
                article["arbitrage_effectue"] = False

    return article


def _demander_arbitrage(titre: str, texte: str, score_llm: int, raison_llm: str,
                          score_estime: float, bloc_few_shot: str = ""):
    """ Second appel : le LLM tranche entre son propre score et l'estimation
        statistique, sur la base des faits — pas en moyennant. """
    extrait = f"Titre : {titre or ''}\n\nTexte : {texte or ''}".strip()[:6000]
    exemples_txt = f"\n{bloc_few_shot}\n" if bloc_few_shot else ""

    systeme = (
        "Tu es un assistant de veille juridique DPO/RGPD. Un premier passage "
        "a donné les résultats suivants pour cet article :\n\n"
        f"- Ton score initial : {score_llm}/10\n"
        f"- Ta raison initiale : {raison_llm}\n"
        f"- Une estimation statistique indépendante (basée sur la similarité "
        f"avec des articles déjà annotés) suggère : {score_estime:.1f}/10\n\n"
        "Ces deux scores divergent significativement. L'estimation statistique "
        "est SOUVENT PEU FIABLE : elle se base sur la ressemblance textuelle "
        "de surface, pas sur la substance réglementaire (elle confond parfois "
        "des sujets proches en thème mais très différents en portée juridique).\n\n"
        "Réexamine l'article et tranche UNIQUEMENT sur la base des faits "
        "suivants, sans chercher à faire un compromis entre les deux scores :\n"
        "- Y a-t-il une décision/sanction d'une autorité identifiée ?\n"
        "- Le texte a-t-il une portée contraignante ou seulement indicative ?\n"
        "- Le lien avec le périmètre DPO est-il direct ou superficiel ?\n"
        f"{exemples_txt}\n"
        "Réponds UNIQUEMENT par un objet JSON valide :\n"
        '{"score_final": <entier 0-10>, "justification_arbitrage": "explique '
        'pourquoi tu confirmes ou corriges ton score initial, en te basant sur '
        'les faits, pas sur l\'estimation statistique"}'
    )

    for tentative in range(1, MAX_RETRIES_LLM + 1):
        time.sleep(DELAI_LLM)
        try:
            resp = openai.chat.completions.create(
                model=LLM_MODELE,
                messages=[
                    {"role": "system", "content": systeme},
                    {"role": "user", "content": extrait},
                ],
                temperature=0,
                response_format={"type": "json_object"},
            )
            contenu = resp.choices[0].message.content.strip()
            verdict = _extraire_json(contenu)
            if verdict is None:
                return None
            try:
                score_final = int(round(float(verdict.get("score_final"))))
            except (TypeError, ValueError):
                return None
            score_final = max(0, min(10, score_final))
            justification = str(verdict.get("justification_arbitrage", "")).strip()
            return score_final, justification
        except openai.RateLimitError:
            time.sleep(DELAI_LLM * tentative * 2)
            continue
        except Exception as e:
            print(f"  ! arbitrage LLM indisponible : {e}")
            return None

    return None


def estimer_score_knn(article_texte, k=10):
    """ Estime le score d'un article par moyenne pondérée des k plus proches
        voisins du dataset annoté, pondérés par leur similarité cosinus.
        Sert de signal statistique complémentaire au jugement du LLM —
        ne remplace pas le score du LLM, vient en renfort/vérification. """
    if _dataset is None:
        construire_index()

    embedder = _get_embedder()
    vecteur_article = embedder.encode(
        [article_texte], normalize_embeddings=True, prompt=INSTRUCTION
    )[0]

    similarites = _embeddings @ vecteur_article
    ordre = np.argsort(-similarites)[:k]

    poids = similarites[ordre]
    poids = np.clip(poids, 0, None)  # ignore les similarités négatives (rare mais possible)

    if poids.sum() == 0:
        return None  # aucun voisin exploitable (tous à similarité nulle/négative)

    scores_voisins = _dataset.iloc[ordre][COL_SCORE].to_numpy(dtype=float)
    score_estime = float(np.average(scores_voisins, weights=poids))

    return {
        "score_estime": round(score_estime, 1),
        "k_utilise": len(ordre),
        "similarite_max": float(similarites[ordre[0]]),
        "similarite_min": float(similarites[ordre[-1]]),
    }



"""
fewshot_dynamique.py — Sélection dynamique d'exemples few-shot par similarité sémantique.

Charge le dataset annoté (CSV/Excel, score 0-10), calcule les embeddings via
Harrier OSS v1 0.6B (sentence-transformers) à partir d'une combinaison
résumé + contenu (en ignorant les placeholders type "pas de résumé
disponible"), les met en cache sur disque, puis permet de récupérer les k
exemples les plus proches d'un article donné pour construire le prompt.
"""
import os
import numpy as np
import pandas as pd
from sentence_transformers import SentenceTransformer

# --- config dataset : à ajuster selon ton fichier ---
DATASET_PATH = "dataset_annote.xlsx"      # ou .csv
COL_TEXTE = "resume"                       # colonne résumé
COL_CONTENU = "contenu"                    # colonne contenu complet
COL_SCORE = "score"                        # colonne score annoté 0-10
COL_TITRE = "titre"                        # colonne titre (pour affichage dans le prompt)

CACHE_PATH = "fewshot_embeddings.npz"      # cache des embeddings calculés

# --- config embedding : Harrier OSS v1 0.6B ---
EMBED_MODEL = "microsoft/harrier-oss-v1-0.6b"
DEVICE = "cuda:1"
INSTRUCTION = (
    "Instruct: Étant donné un article de veille juridique sur la protection "
    "des données personnelles (RGPD/DPO), retrouve les articles annotés les "
    "plus proches en pertinence réglementaire\nQuery: "
)

# valeurs à traiter comme "pas de résumé/contenu" malgré une cellule non vide
PLACEHOLDERS_VIDES = {
    "pas de résumé disponible",
    "pas de resume disponible",
    "non disponible",
    "n/a",
    "",
}

_embedder = None
_dataset = None
_embeddings = None


def _get_embedder():
    global _embedder
    if _embedder is None:
        _embedder = SentenceTransformer(
            EMBED_MODEL,
            device=DEVICE,
            model_kwargs={"dtype": "auto"},
        )
    return _embedder


def _charger_dataset():
    if DATASET_PATH.endswith(".csv"):
        return pd.read_csv(DATASET_PATH)
    return pd.read_excel(DATASET_PATH)


def _texte_valide(valeur) -> str:
    """ Renvoie le texte nettoyé, ou "" si c'est vide/un placeholder connu. """
    t = str(valeur or "").strip()
    if t.lower() in PLACEHOLDERS_VIDES:
        return ""
    return t


def _texte_ligne(ligne):
    """ Texte utilisé pour l'embedding : combinaison résumé + contenu,
        en ignorant les placeholders ('pas de résumé disponible' etc.).
        Fonctionne avec une Series pandas ou un dict. """
    resume = _texte_valide(ligne.get(COL_TEXTE) if hasattr(ligne, "get") else ligne[COL_TEXTE])
    contenu = _texte_valide(ligne.get(COL_CONTENU) if hasattr(ligne, "get") else ligne[COL_CONTENU])
    morceaux = [t for t in (resume, contenu) if t]
    return "\n\n".join(morceaux)


def construire_index(force=False):
    """ Calcule (ou recharge depuis le cache) les embeddings du dataset annoté.
        À appeler une fois au démarrage du script de scoring. """
    global _dataset, _embeddings

    _dataset = _charger_dataset()
    textes = [_texte_ligne(row) for _, row in _dataset.iterrows()]

    if not force and os.path.exists(CACHE_PATH):
        cache = np.load(CACHE_PATH)
        if cache["n"] == len(textes):
            _embeddings = cache["embeddings"]
            return

    embedder = _get_embedder()
    _embeddings = embedder.encode(textes, show_progress_bar=False, normalize_embeddings=True)
    np.savez(CACHE_PATH, embeddings=_embeddings, n=len(textes))


def ajouter_exemple(titre, score, resume=None, contenu=None):
    """ Ajoute un nouvel exemple annoté au dataset + à l'index, sans tout
        recalculer. resume ou contenu peuvent être vides, pas les deux. """
    global _dataset, _embeddings

    if _dataset is None:
        construire_index()

    nouvelle_ligne = {
        COL_TITRE: titre,
        COL_TEXTE: resume or "",
        COL_CONTENU: contenu or "",
        COL_SCORE: score,
    }
    _dataset = pd.concat([_dataset, pd.DataFrame([nouvelle_ligne])], ignore_index=True)

    embedder = _get_embedder()
    texte_pour_embedding = _texte_ligne(nouvelle_ligne)
    nouveau_vecteur = embedder.encode([texte_pour_embedding], normalize_embeddings=True)
    _embeddings = np.vstack([_embeddings, nouveau_vecteur])

    if DATASET_PATH.endswith(".csv"):
        _dataset.to_csv(DATASET_PATH, index=False)
    else:
        _dataset.to_excel(DATASET_PATH, index=False)
    np.savez(CACHE_PATH, embeddings=_embeddings, n=len(_dataset))


def _tranche(score):
    if score <= 3:
        return "0-3"
    if score <= 6:
        return "4-6"
    return "7-10"


def selectionner_exemples(article_texte, k=5, min_par_tranche=1, seuil_similarite=0.0):
    """ Renvoie les k exemples les plus proches de l'article à classer, à
        condition qu'ils dépassent un seuil de similarité minimum. En dessous,
        l'exemple est écarté plutôt que forcé dans le few-shot (peut donc
        renvoyer une liste vide ou plus courte que k).

        seuil_similarite=0.0 par défaut : aucun filtrage réel (le diagnostic
        a montré qu'un seuil strict n'est pas fiable sur ce corpus) — à
        remonter uniquement pour écarter les cas complètement hors-sujet. """
    if _dataset is None:
        construire_index()

    embedder = _get_embedder()
    vecteur_article = embedder.encode(
        [article_texte], normalize_embeddings=True, prompt=INSTRUCTION
    )[0]

    similarites = _embeddings @ vecteur_article  # cosinus (vecteurs déjà normalisés)
    ordre = np.argsort(-similarites)

    resultats = []
    ids_pris = set()
    compte_par_tranche = {}
    for idx in ordre:
        if len(resultats) >= k:
            break
        if similarites[idx] < seuil_similarite:
            break
        ligne = _dataset.iloc[idx]
        resultats.append(ligne)
        ids_pris.add(idx)
        t = _tranche(ligne[COL_SCORE])
        compte_par_tranche[t] = compte_par_tranche.get(t, 0) + 1

    tranches_dataset = _dataset[COL_SCORE].apply(_tranche).unique()
    for t in tranches_dataset:
        if compte_par_tranche.get(t, 0) < min_par_tranche:
            for idx in ordre:
                if idx in ids_pris or similarites[idx] < seuil_similarite:
                    continue
                ligne = _dataset.iloc[idx]
                if _tranche(ligne[COL_SCORE]) == t:
                    resultats.append(ligne)
                    ids_pris.add(idx)
                    compte_par_tranche[t] = compte_par_tranche.get(t, 0) + 1
                    break

    return resultats


def formatter_exemples_pour_prompt(exemples):
    """ Formate les exemples sélectionnés en texte injectable dans le prompt LLM. """
    blocs = []
    for ex in exemples:
        texte = _texte_ligne(ex)
        blocs.append(
            f"Titre : {ex[COL_TITRE]}\n"
            f"Contenu : {texte}\n"
            f"Score attribué : {ex[COL_SCORE]}/10"
        )
    return "\n\n".join(blocs)


import torch
import torch.nn.functional as F
from transformers import AutoTokenizer, AutoModel

EMBED_MODEL = "microsoft/harrier-oss-v1-0.6b"
INSTRUCTION = (
    "Instruct: Étant donné un article de veille juridique sur la protection "
    "des données personnelles (RGPD/DPO), retrouve les articles annotés les "
    "plus proches en pertinence réglementaire\nQuery: "
)

_tokenizer = None
_model = None


def _get_embedder():
    global _tokenizer, _model
    if _model is None:
        _tokenizer = AutoTokenizer.from_pretrained(EMBED_MODEL)
        _model = AutoModel.from_pretrained(EMBED_MODEL, dtype="auto")
        _model.eval()
        if torch.cuda.is_available():
            _model.cuda()
    return _tokenizer, _model


def _last_token_pool(last_hidden_state, attention_mask):
    left_padding = attention_mask[:, -1].sum() == attention_mask.shape[0]
    if left_padding:
        return last_hidden_state[:, -1]
    sequence_lengths = attention_mask.sum(dim=1) - 1
    batch_size = last_hidden_state.shape[0]
    return last_hidden_state[torch.arange(batch_size), sequence_lengths]


def _encoder(textes, avec_instruction=False):
    """ avec_instruction=True pour l'article à classer (query),
        False pour les documents du dataset annoté. """
    tokenizer, model = _get_embedder()
    if avec_instruction:
        textes = [INSTRUCTION + t for t in textes]

    batch = tokenizer(textes, max_length=8192, padding=True, truncation=True, return_tensors="pt")
    if torch.cuda.is_available():
        batch = {k: v.cuda() for k, v in batch.items()}

    with torch.no_grad():
        outputs = model(**batch)
        embeddings = _last_token_pool(outputs.last_hidden_state, batch["attention_mask"])
        embeddings = F.normalize(embeddings, p=2, dim=1)

    return embeddings.cpu().numpy()

let out="";
        if(isAbsent&&detail.length){
          // Regrouper par contributeur : un sous-bloc par personne
          const parContrib={};
          const sansContrib=[];
          detail.forEach(p=>{
            const cs=String(p.contributeurs_sujet||"").split(";")
                       .map(s=>s.trim()).filter(Boolean)
                       .filter(c=>c.toLowerCase()!==String(sel).trim().toLowerCase());
            if(cs.length){cs.forEach(c=>{(parContrib[c]=parContrib[c]||[]).push(p);});}
            else sansContrib.push(p);
          });
          Object.keys(parContrib).sort().forEach(c=>{
            out+=sousTitre(`renseignés par ${esc(c)} (${parContrib[c].length})`,"var(--cyan)");
            out+=parContrib[c].map(p=>projItemFn(p,"")).join("");
          });
          if(sansContrib.length){
            out+=sousTitre(`renseignés (${sansContrib.length})`,"var(--cyan)");
            out+=sansContrib.map(p=>projItemFn(p,"")).join("");
          }
        } else {
          out+=detail.map(p=>projItemFn(p,"")).join("");
        }
__
# Contributeurs de CE sujet : qui a saisi une ligne pour lui
        _contribs = []
        if "source_fichier" in groupe.columns:
            for src in groupe["source_fichier"].dropna().unique():
                a = re.sub(r"(?i)^fiches?_monito", "", str(src)).strip()
                a = re.sub(r"(?i)[.]xls[xm]?$", "", a)
                a = a.replace("_", " ").strip()
                if a and a.lower() != "template":
                    _contribs.append(a)
        row["contributeurs_sujet"] = "; ".join(sorted(set(_contribs)))



// ── Personnes présentes uniquement dans l'historique (aucun sujet cette quinzaine) ──
  // (celles qui ont déjà une carte normale sont exclues → une seule carte par personne)
  const _nomsAvecCarte=new Set(resps.map(([name])=>name));
  const absentsSansCarte=Object.entries(respHistorique)
    .filter(([name])=>!_nomsAvecCarte.has(name))
    .sort();
____

<div class="collab-grid">
      ${resps.map(([name,r])=>{
        const abs=!_aRempli(name);
        return `
        <div class="collab-card ${name===sel?"selected":""}" onclick="renderCollabs('${esc(name)}')"
             ${abs?'style="border-style:dashed"':""}>
          <div class="collab-header">
            <div class="avatar" style="background:${respColor(name)}22;color:${respColor(name)}">${initials(name)}</div>
            <div>
              <div class="collab-name">${esc(name)}</div>
              <div class="collab-sub">${r.total} sujet${r.total>1?"s":""} · ${r.en_cours||0} actif${(r.en_cours||0)>1?"s":""}</div>
              ${abs?`<div style="font-size:9px;color:var(--amber);font-family:var(--font-mono);margin-top:2px">⚠ n'a pas rempli cette quinzaine</div>`:""}
            </div>
          </div>
          <div style="display:flex;gap:4px;flex-wrap:wrap;margin-bottom:6px">
            ${(r.domaines||[]).slice(0,3).map(d=>`<span style="font-size:9px;padding:2px 6px;border-radius:10px;background:var(--bg4);color:var(--text3);font-family:var(--font-mono)">${esc(d)}</span>`).join("")}
          </div>
          <div class="charge-bar"><div class="charge-fill" style="width:${Math.round((r.en_cours||0)/maxE*100)}%;background:${abs?"var(--amber)":respColor(name)}"></div></div>
        </div>`;}).join("")}
      ${absentsSansCarte.map(([name,lastRow])=>`
        <div class="collab-card ${name===sel?"selected":""}" onclick="renderCollabs('${esc(name)}')"
             style="opacity:.65;border-style:dashed">
          <div class="collab-header">
            <div class="avatar" style="background:${respColor(name)}22;color:${respColor(name)}">${initials(name)}</div>
            <div>
              <div class="collab-name">${esc(name)}</div>
              <div class="collab-sub" style="color:var(--amber)">Absent cette quinzaine</div>
              <div style="font-size:8px;color:var(--text3);font-family:var(--font-mono)">
                dernier suivi : ${esc(lastRow.quinzaine)}
              </div>
            </div>
          </div>
          <div class="charge-bar"><div class="charge-fill" style="width:0%;background:var(--amber)"></div></div>
        </div>`).join("")}
    </div>



// ── Détail du collaborateur sélectionné ─────────────────────────
  const _tri=(a,b)=>({"En retard":0,"À risque":1,"En cours":2,"Stand by":3,"Terminé":4}[a.statut]||9)-
                    ({"En retard":0,"À risque":1,"En cours":2,"Stand by":3,"Terminé":4}[b.statut]||9);

  // Bloc 1 : sujets renseignés cette quinzaine (par lui ou par un collaborateur)
  let detail=DATA.projets.filter(p=>p.responsable_principal===sel).sort(_tri);

  // Absent = n'a rempli aucune fiche cette quinzaine
  const isAbsent=!_aRempli(sel);

  // Bloc 2 : ses autres sujets, repris de la dernière quinzaine qu'il a lui-même remplie
  let detailHist=[];
  let lastQAbs="";
  if(isAbsent){
    (DATA.quinzaines||[]).filter(q=>q<quinzaineActive).sort().forEach(q=>{
      const snap=DATA.snapshots[q];
      const auteurs=new Set();
      (snap?.projets||[]).forEach(p=>{
        String(p.auteurs_quinzaine||"").split(";").map(s=>s.trim().toLowerCase())
          .filter(Boolean).forEach(a=>auteurs.add(a));
      });
      if(auteurs.has(String(sel).trim().toLowerCase()))lastQAbs=q;
    });
    if(lastQAbs){
      const dejaVus=new Set(detail.map(p=>p.projet_id||p.ref_sujet));
      detailHist=(DATA.snapshots[lastQAbs]?.projets||[])
        .filter(p=>p.responsable_principal===sel)
        .filter(p=>!dejaVus.has(p.projet_id||p.ref_sujet))   // pas de doublon avec le bloc 1
        .sort(_tri);
    }
  }

_____
<div class="card">
      <div class="card-title">sujets :: ${esc(sel)} (${detail.length+detailHist.length})</div>
      ${isAbsent?`<div style="font-size:10px;color:var(--amber);font-family:var(--font-mono);margin-bottom:10px;padding:6px 10px;background:var(--amber-dim);border-radius:var(--radius);border:1px solid rgba(245,158,11,.2)">
        ⚠ ${esc(sel)} n'a rempli aucune fiche cette quinzaine
      </div>`:""}
      <div class="proj-list">
        ${(()=>{
        const projItemFn=(p,histQ)=>{
          const partage=(p.partage_prochain_point||"").toString().toLowerCase().trim();
          const aPartager=partage==="oui"||partage==="yes"||partage==="1";
          let html=projItem(p);
          const id=p.projet_id||p.ref_sujet;
          if(histQ){
            html=html.replace(`openModal('${esc(id)}')`,`openModalHistorique('${esc(id)}','${esc(histQ)}')`);
          }
          if(aPartager){
            html=html.replace('class="proj-item"','class="proj-item proj-partage"');
            const call=histQ?`openModalHistorique('${esc(id)}','${esc(histQ)}')`:`openModal('${esc(id)}')`;
            html=html.replace(`onclick="${call}"`,`onclick="marquerTraite(this);${call}"`);
            html=html.replace('</div>',`<span style="font-size:9px;background:var(--amber-dim);color:var(--amber);border:1px solid rgba(245,158,11,.3);padding:1px 5px;border-radius:8px;font-family:var(--font-mono);flex-shrink:0;margin-left:auto">!!! à partager</span></div>`);
          }
          return html;
        };
        const sousTitre=(txt,coul)=>`<div style="font-size:9px;font-weight:600;text-transform:uppercase;letter-spacing:.08em;color:${coul};font-family:var(--font-mono);padding:10px 0 5px">${txt}</div>`;
        let out="";
        if(isAbsent&&detail.length){
          out+=sousTitre(`renseignés par un collaborateur (${detail.length})`,"var(--cyan)");
        }
        out+=detail.map(p=>projItemFn(p,"")).join("");
        if(detailHist.length){
          out+=sousTitre(`dernière quinzaine remplie — ${esc(lastQAbs)} (${detailHist.length})`,"var(--amber)");
          out+=detailHist.map(p=>projItemFn(p,lastQAbs)).join("");
        }
        return out||'<div style="color:var(--text3);font-size:11px;font-family:var(--font-mono);padding:8px">// aucun sujet</div>';
        })()}
      </div>
    </div>`;

___


// Absent = n'a rempli aucune fiche cette quinzaine (même si ses sujets
  // ont été renseignés par un collaborateur)
  let isAbsent=!_aRempli(sel);
  if(isAbsent){
    // Retrouver la dernière quinzaine où cette personne a réellement rempli
    let lastQ="";
    (DATA.quinzaines||[]).filter(q=>q<quinzaineActive).sort().forEach(q=>{
      const snap=DATA.snapshots[q];
      const auteurs=new Set();
      (snap?.projets||[]).forEach(p=>{
        String(p.auteurs_quinzaine||"").split(";").map(s=>s.trim().toLowerCase())
          .filter(Boolean).forEach(a=>auteurs.add(a));
      });
      if(auteurs.has(String(sel).trim().toLowerCase()))lastQ=q;
    });
    if(lastQ){
      const snap=DATA.snapshots[lastQ];
      detail=(snap?.projets||[]).filter(p=>p.responsable_principal===sel)
        .sort((a,b)=>({"En retard":0,"À risque":1,"En cours":2,"Stand by":3,"Terminé":4}[a.statut]||9)-
                      ({"En retard":0,"À risque":1,"En cours":2,"Stand by":3,"Terminé":4}[b.statut]||9));
    }
}


# ── Auteurs ayant réellement rempli une fiche pour cette quinzaine ──
        # (avant fusion : dfs contient encore le source_fichier de chaque ligne)
        auteurs_q = set()
        for _df in dfs:
            if "source_fichier" in _df.columns:
                for src in _df["source_fichier"].dropna().unique():
                    a = re.sub(r"(?i)^fiches?_monito", "", str(src)).strip()
                    a = re.sub(r"(?i)[.]xls[xm]?$", "", a)
                    a = a.replace("_", " ").strip()
                    if a:
                        auteurs_q.add(a)
        df_c["auteurs_quinzaine"] = "; ".join(sorted(auteurs_q))

________

# Le référentiel fait autorité sur le responsable principal
            if "responsable_principal_meta" in df_c.columns:
                mask = df_c["responsable_principal_meta"].notna() & \
                       (df_c["responsable_principal_meta"].astype(str).str.strip() != "")
                df_c.loc[mask, "responsable_principal"] = df_c.loc[mask, "responsable_principal_meta"]

___

// ── Collaborateurs absents (n'ont rempli aucune fiche cette quinzaine) ──
  // Source : auteurs_quinzaine, produit par le parser depuis source_fichier
  const _auteursQ=new Set();
  (DATA.projets||[]).forEach(p=>{
    String(p.auteurs_quinzaine||"").split(";").map(s=>s.trim().toLowerCase())
      .filter(Boolean).forEach(a=>_auteursQ.add(a));
  });
  const _aRempli=nom=>{
    const n=String(nom||"").trim().toLowerCase();
    if(!n)return false;
    if(_auteursQ.size===0)return true;   // pas d'info → ne marquer personne absent
    return _auteursQ.has(n);
  };
  const respActifs=new Set(
    DATA.projets.map(p=>p.responsable_principal).filter(Boolean).filter(_aRempli)
  );

"""
fewshot_dynamique.py — Sélection dynamique d'exemples few-shot par similarité sémantique.

Charge le dataset annoté (CSV/Excel, score 0-10), calcule les embeddings à
partir d'une combinaison résumé + contenu (en ignorant les placeholders type
"pas de résumé disponible"), les met en cache sur disque, puis permet de
récupérer les k exemples les plus proches d'un article donné pour construire
le prompt.
"""
import os
import numpy as np
import pandas as pd
from sentence_transformers import SentenceTransformer

# --- config : à ajuster selon ton fichier ---
DATASET_PATH = "dataset_annote.xlsx"      # ou .csv
COL_TEXTE = "resume"                       # colonne résumé
COL_CONTENU = "contenu"                    # colonne contenu complet
COL_SCORE = "score"                        # colonne score annoté 0-10
COL_TITRE = "titre"                        # colonne titre (pour affichage dans le prompt)

CACHE_PATH = "fewshot_embeddings.npz"      # cache des embeddings calculés
EMBED_MODEL = "paraphrase-multilingual-MiniLM-L12-v2"  # léger, multilingue, tourne en local

# valeurs à traiter comme "pas de résumé/contenu" malgré une cellule non vide
PLACEHOLDERS_VIDES = {
    "pas de résumé disponible",
    "pas de resume disponible",
    "non disponible",
    "n/a",
    "",
}

_embedder = None
_dataset = None
_embeddings = None


def _get_embedder():
    global _embedder
    if _embedder is None:
        _embedder = SentenceTransformer(EMBED_MODEL)
    return _embedder


def _charger_dataset():
    if DATASET_PATH.endswith(".csv"):
        return pd.read_csv(DATASET_PATH)
    return pd.read_excel(DATASET_PATH)


def _texte_valide(valeur) -> str:
    """ Renvoie le texte nettoyé, ou "" si c'est vide/un placeholder connu. """
    t = str(valeur or "").strip()
    if t.lower() in PLACEHOLDERS_VIDES:
        return ""
    return t


def _texte_ligne(ligne):
    """ Texte utilisé pour l'embedding : combinaison résumé + contenu,
        en ignorant les placeholders ('pas de résumé disponible' etc.).
        Fonctionne avec une Series pandas ou un dict. """
    resume = _texte_valide(ligne.get(COL_TEXTE) if hasattr(ligne, "get") else ligne[COL_TEXTE])
    contenu = _texte_valide(ligne.get(COL_CONTENU) if hasattr(ligne, "get") else ligne[COL_CONTENU])
    morceaux = [t for t in (resume, contenu) if t]
    return "\n\n".join(morceaux)


def construire_index(force=False):
    """ Calcule (ou recharge depuis le cache) les embeddings du dataset annoté.
        À appeler une fois au démarrage du script de scoring. """
    global _dataset, _embeddings

    _dataset = _charger_dataset()
    textes = [_texte_ligne(row) for _, row in _dataset.iterrows()]

    if not force and os.path.exists(CACHE_PATH):
        cache = np.load(CACHE_PATH)
        if cache["n"] == len(textes):
            _embeddings = cache["embeddings"]
            return

    embedder = _get_embedder()
    _embeddings = embedder.encode(textes, show_progress_bar=False, normalize_embeddings=True)
    np.savez(CACHE_PATH, embeddings=_embeddings, n=len(textes))


def ajouter_exemple(titre, score, resume=None, contenu=None):
    """ Ajoute un nouvel exemple annoté au dataset + à l'index, sans tout
        recalculer. resume ou contenu peuvent être vides, pas les deux. """
    global _dataset, _embeddings

    if _dataset is None:
        construire_index()

    nouvelle_ligne = {
        COL_TITRE: titre,
        COL_TEXTE: resume or "",
        COL_CONTENU: contenu or "",
        COL_SCORE: score,
    }
    _dataset = pd.concat([_dataset, pd.DataFrame([nouvelle_ligne])], ignore_index=True)

    embedder = _get_embedder()
    texte_pour_embedding = _texte_ligne(nouvelle_ligne)
    nouveau_vecteur = embedder.encode([texte_pour_embedding], normalize_embeddings=True)
    _embeddings = np.vstack([_embeddings, nouveau_vecteur])

    if DATASET_PATH.endswith(".csv"):
        _dataset.to_csv(DATASET_PATH, index=False)
    else:
        _dataset.to_excel(DATASET_PATH, index=False)
    np.savez(CACHE_PATH, embeddings=_embeddings, n=len(_dataset))


def _tranche(score):
    if score <= 3:
        return "0-3"
    if score <= 6:
        return "4-6"
    return "7-10"


def selectionner_exemples(article_texte, k=5, min_par_tranche=1, seuil_similarite=0.0):
    """ Renvoie les k exemples les plus proches de l'article à classer, à
        condition qu'ils dépassent un seuil de similarité minimum. En dessous,
        l'exemple est écarté plutôt que forcé dans le few-shot (peut donc
        renvoyer une liste vide ou plus courte que k).

        seuil_similarite=0.0 par défaut : aucun filtrage réel (le diagnostic
        a montré qu'un seuil strict n'est pas fiable sur ce corpus) — à
        remonter uniquement pour écarter les cas complètement hors-sujet. """
    if _dataset is None:
        construire_index()

    embedder = _get_embedder()
    vecteur_article = embedder.encode([article_texte], normalize_embeddings=True)[0]

    similarites = _embeddings @ vecteur_article  # cosinus (vecteurs déjà normalisés)
    ordre = np.argsort(-similarites)

    resultats = []
    ids_pris = set()
    compte_par_tranche = {}
    for idx in ordre:
        if len(resultats) >= k:
            break
        if similarites[idx] < seuil_similarite:
            break  # au-delà, plus rien n'est assez proche — on arrête
        ligne = _dataset.iloc[idx]
        resultats.append(ligne)
        ids_pris.add(idx)
        t = _tranche(ligne[COL_SCORE])
        compte_par_tranche[t] = compte_par_tranche.get(t, 0) + 1

    # rattrapage min_par_tranche, avec le même seuil de similarité
    tranches_dataset = _dataset[COL_SCORE].apply(_tranche).unique()
    for t in tranches_dataset:
        if compte_par_tranche.get(t, 0) < min_par_tranche:
            for idx in ordre:
                if idx in ids_pris or similarites[idx] < seuil_similarite:
                    continue
                ligne = _dataset.iloc[idx]
                if _tranche(ligne[COL_SCORE]) == t:
                    resultats.append(ligne)
                    ids_pris.add(idx)
                    compte_par_tranche[t] = compte_par_tranche.get(t, 0) + 1
                    break

    return resultats


def formatter_exemples_pour_prompt(exemples):
    """ Formate les exemples sélectionnés en texte injectable dans le prompt LLM. """
    blocs = []
    for ex in exemples:
        texte = _texte_ligne(ex)
        blocs.append(
            f"Titre : {ex[COL_TITRE]}\n"
            f"Contenu : {texte}\n"
            f"Score attribué : {ex[COL_SCORE]}/10"
        )
    return "\n\n".join(blocs)


"""
fewshot_dynamique.py — Sélection dynamique d'exemples few-shot par similarité sémantique.

Charge le dataset annoté (CSV/Excel, score 0-10), calcule les embeddings une
fois, les met en cache sur disque, puis permet de récupérer les k exemples
les plus proches d'un article donné pour construire le prompt.
"""
import os
import numpy as np
import pandas as pd
from sentence_transformers import SentenceTransformer

# --- config : à ajuster selon ton fichier ---
DATASET_PATH = "dataset_annote.xlsx"      # ou .csv
COL_TEXTE = "resume"                       # colonne texte utilisée pour le matching
COL_SCORE = "score"                        # colonne score annoté 0-10
COL_TITRE = "titre"                        # colonne titre (pour affichage dans le prompt)

CACHE_PATH = "fewshot_embeddings.npz"      # cache des embeddings calculés
EMBED_MODEL = "paraphrase-multilingual-MiniLM-L12-v2"  # léger, multilingue, tourne en local

_embedder = None
_dataset = None
_embeddings = None


def _get_embedder():
    global _embedder
    if _embedder is None:
        _embedder = SentenceTransformer(EMBED_MODEL)
    return _embedder


def _charger_dataset():
    if DATASET_PATH.endswith(".csv"):
        return pd.read_csv(DATASET_PATH)
    return pd.read_excel(DATASET_PATH)


def construire_index(force=False):
    """ Calcule (ou recharge depuis le cache) les embeddings du dataset annoté.
        À appeler une fois au démarrage du script de scoring. """
    global _dataset, _embeddings

    _dataset = _charger_dataset()
    textes = _dataset[COL_TEXTE].fillna("").astype(str).tolist()

    if not force and os.path.exists(CACHE_PATH):
        cache = np.load(CACHE_PATH)
        if cache["n"] == len(textes):
            _embeddings = cache["embeddings"]
            return

    embedder = _get_embedder()
    _embeddings = embedder.encode(textes, show_progress_bar=False, normalize_embeddings=True)
    np.savez(CACHE_PATH, embeddings=_embeddings, n=len(textes))


def ajouter_exemple(titre, resume, score):
    """ Ajoute un nouvel exemple annoté au dataset + à l'index, sans tout
        recalculer. Pratique pour enrichir le dataset au fil de l'eau. """
    global _dataset, _embeddings

    if _dataset is None:
        construire_index()

    nouvelle_ligne = {COL_TITRE: titre, COL_TEXTE: resume, COL_SCORE: score}
    _dataset = pd.concat([_dataset, pd.DataFrame([nouvelle_ligne])], ignore_index=True)

    embedder = _get_embedder()
    nouveau_vecteur = embedder.encode([resume], normalize_embeddings=True)
    _embeddings = np.vstack([_embeddings, nouveau_vecteur])

    if DATASET_PATH.endswith(".csv"):
        _dataset.to_csv(DATASET_PATH, index=False)
    else:
        _dataset.to_excel(DATASET_PATH, index=False)
    np.savez(CACHE_PATH, embeddings=_embeddings, n=len(_dataset))


def _tranche(score):
    if score <= 3:
        return "0-3"
    if score <= 6:
        return "4-6"
    return "7-10"


def selectionner_exemples(article_texte, k=5, min_par_tranche=1, seuil_similarite=0.35):
    """ Renvoie les k exemples les plus proches de l'article à classer, à
        condition qu'ils dépassent un seuil de similarité minimum. En dessous,
        l'exemple est écarté plutôt que forcé dans le few-shot (peut donc
        renvoyer une liste vide ou plus courte que k). """
    if _dataset is None:
        construire_index()

    embedder = _get_embedder()
    vecteur_article = embedder.encode([article_texte], normalize_embeddings=True)[0]

    similarites = _embeddings @ vecteur_article  # cosinus (vecteurs déjà normalisés)
    ordre = np.argsort(-similarites)

    resultats = []
    ids_pris = set()
    compte_par_tranche = {}
    for idx in ordre:
        if len(resultats) >= k:
            break
        if similarites[idx] < seuil_similarite:
            break  # au-delà, plus rien n'est assez proche — on arrête
        ligne = _dataset.iloc[idx]
        resultats.append(ligne)
        ids_pris.add(idx)
        t = _tranche(ligne[COL_SCORE])
        compte_par_tranche[t] = compte_par_tranche.get(t, 0) + 1

    # rattrapage min_par_tranche, avec le même seuil de similarité
    tranches_dataset = _dataset[COL_SCORE].apply(_tranche).unique()
    for t in tranches_dataset:
        if compte_par_tranche.get(t, 0) < min_par_tranche:
            for idx in ordre:
                if idx in ids_pris or similarites[idx] < seuil_similarite:
                    continue
                ligne = _dataset.iloc[idx]
                if _tranche(ligne[COL_SCORE]) == t:
                    resultats.append(ligne)
                    ids_pris.add(idx)
                    compte_par_tranche[t] = compte_par_tranche.get(t, 0) + 1
                    break

    return resultats


def formatter_exemples_pour_prompt(exemples):
    """ Formate les exemples sélectionnés en texte injectable dans le prompt LLM. """
    blocs = []
    for ex in exemples:
        blocs.append(
            f"Titre : {ex[COL_TITRE]}\n"
            f"Résumé : {ex[COL_TEXTE]}\n"
            f"Score attribué : {ex[COL_SCORE]}/10"
        )
    return "\n\n".join(blocs)






---------------------------------------



"""
diagnostic_similarite.py — Analyse la distribution des similarités sémantiques
du dataset annoté pour choisir un seuil de similarité justifié, plutôt qu'estimé.

Usage : python diagnostic_similarite.py
"""
import numpy as np
import fewshot_dynamique as fs

fs.construire_index(force=True)

dataset = fs._dataset
embeddings = fs._embeddings
scores = dataset[fs.COL_SCORE].to_numpy()

n = len(dataset)
print(f"Dataset : {n} articles annotés\n")

# --- matrice de similarité complète (toutes les paires i < j) ---
sim_matrix = embeddings @ embeddings.T

paires_sim = []
paires_ecart = []
for i in range(n):
    for j in range(i + 1, n):
        paires_sim.append(sim_matrix[i, j])
        paires_ecart.append(abs(scores[i] - scores[j]))

paires_sim = np.array(paires_sim)
paires_ecart = np.array(paires_ecart)

# --- regroupement par écart de score ---
groupes = {
    "écart faible (0-2)":  paires_sim[paires_ecart <= 2],
    "écart moyen (3-5)":   paires_sim[(paires_ecart > 2) & (paires_ecart <= 5)],
    "écart fort (6-10)":   paires_sim[paires_ecart > 5],
}

print(f"{'Groupe':<22} {'n paires':>9} {'moyenne':>9} {'médiane':>9} {'p25':>7} {'p75':>7} {'p90':>7}")
for nom, valeurs in groupes.items():
    if len(valeurs) == 0:
        print(f"{nom:<22} {'—':>9}")
        continue
    print(f"{nom:<22} {len(valeurs):>9} "
          f"{valeurs.mean():>9.3f} {np.median(valeurs):>9.3f} "
          f"{np.percentile(valeurs, 25):>7.3f} {np.percentile(valeurs, 75):>7.3f} "
          f"{np.percentile(valeurs, 90):>7.3f}")

# --- suggestion de seuil ---
ecart_fort = groupes["écart fort (6-10)"]
ecart_faible = groupes["écart faible (0-2)"]

if len(ecart_fort) > 0 and len(ecart_faible) > 0:
    seuil_suggere = np.percentile(ecart_fort, 90)
    recouvrement = (ecart_faible < seuil_suggere).mean() * 100
    print(f"\nSeuil suggéré : {seuil_suggere:.3f}")
    print(f"→ Au-dessus de ce seuil, seuls ~10% des paires à écart fort passeraient encore.")
    print(f"→ Mais {recouvrement:.0f}% des paires à écart faible seraient EXCLUES par ce seuil.")
    if recouvrement > 30:
        print("  (recouvrement élevé : les deux groupes se chevauchent beaucoup, "
              "les embeddings séparent mal la pertinence par le score sur ce corpus — "
              "seuil à interpréter avec prudence, ou revoir la stratégie.)")
else:
    print("\nPas assez de paires dans un groupe pour suggérer un seuil "
          "(dataset trop petit ou peu de variance de score).")



------—--------------







from pptx.enum.text import MSO_ANCHOR

def remplir_article(slide, art, trimestre=None):
    t = _shape_by_name(slide, "Titre 6")
    if t:
        _set_text(t, art.get("titre") or "Sans titre")
    corps = _shape_by_name(slide, "Rectangle 3")
    if corps:
        corps.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        corps.text_frame.margin_top = Pt(12)
        objet = _extraire_objet(art.get("llm_resume") or "")
        _set_text(corps, objet or art.get("resume") or "")


def _set_text(shape, text):
    """ Remplace le texte d'une shape en gardant le style du 1er run.

    Gère le multi-lignes : chaque '\n' devient un nouveau paragraphe,
    en réutilisant le style du paragraphe modèle (y compris pour les
    paragraphes ajoutés, dont la mise en forme est clonée sur le modèle).
    """
    tf = shape.text_frame
    lignes = (text or "").split("\n")

    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = lignes[0]
        for r in p0.runs[1:]:
            r.text = ""
    else:
        p0.add_run().text = lignes[0]

    # on garde une référence au XML du paragraphe modèle AVANT de supprimer les autres
    modele_pPr = copy.deepcopy(p0._p.find(qn("a:pPr"))) if p0._p.find(qn("a:pPr")) is not None else None
    modele_rPr = copy.deepcopy(p0.runs[0]._r.find(qn("a:rPr"))) if p0.runs and p0.runs[0]._r.find(qn("a:rPr")) is not None else None

    # supprimer les paragraphes existants au-delà du premier
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)

    # ajouter un paragraphe par ligne supplémentaire, en clonant le style du modèle
    for ligne in lignes[1:]:
        p = tf.add_paragraph()
        if modele_pPr is not None:
            old_pPr = p._p.find(qn("a:pPr"))
            if old_pPr is not None:
                p._p.remove(old_pPr)
            p._p.insert(0, copy.deepcopy(modele_pPr))
        r = p.add_run()
        r.text = ligne
        if modele_rPr is not None:
            r._r.insert(0, copy.deepcopy(modele_rPr))


_LABELS = ("Titre", "Type", "Objet", "Portée")

def _extraire_champ(llm_resume: str, nom_champ: str) -> str:
    """ Extrait le contenu d'un champ 'Nom : ...' du résumé structuré,
        même sur plusieurs lignes/paragraphes, jusqu'au prochain label connu. """
    if not llm_resume:
        return ""
    autres = "|".join(l for l in _LABELS if l != nom_champ)
    pattern = rf"(?im)^\s*[-•*]?\s*{nom_champ}\s*:\s*(.+?)(?=\n\s*[-•*]?\s*(?:{autres})\s*:|\Z)"
    m = re.search(pattern, llm_resume, flags=re.DOTALL)
    if m:
        return m.group(1).strip()
    return ""


def _extraire_objet(llm_resume: str) -> str:
    champ = _extraire_champ(llm_resume, "Objet")
    return champ or llm_resume.strip()


def _extraire_titre(llm_resume: str) -> str:
    return _extraire_champ(llm_resume, "Titre")



def _plage_trimestre(trimestre: str):
    """ '2026 – 2e Trimestre' -> ('2026-04-01', '2026-06-30'). None si illisible. """
    if not trimestre:
        return None
    m_an = re.search(r"(20\d{2})", trimestre)
    m_tr = re.search(r"([1-4])\s*(?:er|e|ème|eme)?\s*Trimestre", trimestre, re.I)
    if not (m_an and m_tr):
        return None
    an, tr = int(m_an.group(1)), int(m_tr.group(1))
    debuts = {1: "01-01", 2: "04-01", 3: "07-01", 4: "10-01"}
    fins   = {1: "03-31", 2: "06-30", 3: "09-30", 4: "12-31"}
    return f"{an}-{debuts[tr]}", f"{an}-{fins[tr]}"

def collecter(top, seuil, trimestre=None):
    plage = _plage_trimestre(trimestre)
    conn = get_connection()
    try:
        sql = """
            SELECT a.*, s.nom AS source_nom,
                   substr(a.date_publication,7,4) || '-' ||
                   substr(a.date_publication,1,2) || '-' ||
                   substr(a.date_publication,4,2) AS date_iso
            FROM articles a JOIN sources s ON s.id = a.source_id
        """
        params = ()
        if plage:
            sql += " WHERE date_iso BETWEEN ? AND ?"
            params = plage
        rows = [dict(r) for r in conn.execute(sql, params)]
    finally:
        conn.close()
    # ... la suite inchangée





def _traiter_articles(articles, libelle):
    print(f"{len(articles)} articles à {libelle}")
    for i, article in enumerate(articles, 1):
        article = enrichir_contenu(article)
        update_contenu(article)
        article = scorer_article(article)

        # on mémorise le résumé existant pour ne pas l'écraser
        ancien_resume = article.get("llm_resume")

        article = scorer_llm_article(article)      # recalcule le score

        # si un résumé existait déjà, on le restaure ; sinon on garde le nouveau
        if ancien_resume:
            article["llm_resume"] = ancien_resume

        update_scores(article)
        if i % 10 == 0:
            print(f"  ... {i}/{len(articles)}")
    print(f"✓ {len(articles)} articles traités")


from pptx.dml.color import RGBColor

def remplir_sommaire(slide, articles, slides_articles, prs):
    """ Remplit le sommaire ET crée les liens cliquables vers chaque article. """
    contenu = _shape_by_name(slide, "Espace réservé du contenu 4")
    if not contenu:
        return
    tf = contenu.text_frame
    modele = tf.paragraphs[0]
    style_run = modele.runs[0] if modele.runs else None

    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)

    if not articles:
        if style_run:
            style_run.text = "—"
        return

    def _style_lien(run):
        # police / couleur des liens du sommaire — à ajuster ici
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # ex. noir
        run.font.underline = False
        # run.font.name = "Calibri"
        # run.font.size = Pt(14)
        # run.font.bold = False

    # 1ère ligne : réutilise le paragraphe modèle
    if style_run:
        style_run.text = articles[0].get("titre", "")[:90]
        _lien_vers_slide(style_run, slides_articles[0], prs)
        _style_lien(style_run)

    # lignes suivantes
    for art, sl in zip(articles[1:], slides_articles[1:]):
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = art.get("titre", "")[:90]
        _lien_vers_slide(r, sl, prs)
        _style_lien(r)
      
  """
generer_rapport_template.py — Remplit le template Cofidis avec les articles de la base.

Structure du template (11 slides) :
  1  couverture
  2  séparateur Partie 01 (FR)      3 sommaire      4 article
  5  séparateur Partie 02 (Intl)    6 sommaire      7 article
  8  séparateur Partie 03 (Juri)    9 sommaire     10 article
  11 contacts

Pour chaque partie : on garde le séparateur + le sommaire, et on duplique
la slide "article" autant de fois qu'il y a d'articles retenus.

Usage : python generer_rapport_template.py [--top N] [--seuil S] [--trimestre "2026 – 2e Trimestre"]
"""
import copy
import re
import sys

from pptx import Presentation
from pptx.util import Pt

from database import get_connection
from classer_rapport import classer

TEMPLATE = "template.pptx"
SORTIE = "rapport_veille_cofidis.pptx"

# index (0-based) des slides du template
SEP = {1: 1, 2: 4, 3: 7}       # séparateur de chaque partie
SOMMAIRE = {1: 2, 2: 5, 3: 8}  # sommaire de chaque partie
ARTICLE = {1: 3, 2: 6, 3: 9}   # slide article-type de chaque partie


def meilleur_score(a):
    return a["llm_score"] if a.get("llm_score") is not None else (a.get("score_pertinence") or 0)


def collecter(top, seuil):
    conn = get_connection()
    try:
        rows = [dict(r) for r in conn.execute(
            "SELECT a.*, s.nom AS source_nom FROM articles a JOIN sources s ON s.id=a.source_id")]
    finally:
        conn.close()
    retenus = [a for a in rows if
               (a.get("llm_score") is not None and a["llm_score"] >= seuil)
               or (a.get("llm_score") is None and (a.get("score_pertinence") or 0) > 0)]
    parties = {1: [], 2: [], 3: []}
    for a in retenus:
        parties[classer(a)].append(a)
    for p in parties:
        parties[p].sort(key=meilleur_score, reverse=True)
        parties[p] = parties[p][:top]
    return parties


# ---- helpers python-pptx pour dupliquer/supprimer des slides ----

def _dup_slide(prs, index):
    """ Duplique la slide à l'index donné, l'ajoute en fin, renvoie la nouvelle. """
    source = prs.slides[index]
    blank = source.slide_layout
    new = prs.slides.add_slide(blank)
    # vider les placeholders ajoutés par le layout
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    # copier toutes les shapes de la source
    for sh in source.shapes:
        new.shapes._spTree.append(copy.deepcopy(sh._element))
    return new


def _xml_move(prs, from_pos, to_pos):
    """ Déplace la slide de from_pos vers to_pos dans l'ordre. """
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    el = ids[from_pos]
    sldIdLst.remove(el)
    sldIdLst.insert(to_pos, el)


def _set_text(shape, text):
    """ Remplace le texte d'une shape en gardant le style du 1er run.

    Gère le multi-lignes : chaque '\n' devient un nouveau paragraphe,
    en réutilisant le style du paragraphe modèle.
    """
    tf = shape.text_frame
    lignes = (text or "").split("\n")

    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = lignes[0]
        for r in p0.runs[1:]:
            r.text = ""
    else:
        p0.add_run().text = lignes[0]

    # supprimer les paragraphes existants au-delà du premier
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)

    # ajouter un paragraphe par ligne supplémentaire
    for ligne in lignes[1:]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = ligne


def _shape_by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def _extraire_objet(llm_resume: str) -> str:
    """ Extrait la ligne 'Objet : ...' du résumé structuré.
        Renvoie le résumé entier si le champ Objet est absent. """
    if not llm_resume:
        return ""
    for ligne in llm_resume.splitlines():
        l = ligne.strip().lstrip("-•* ").strip()
        # tolère "Objet :", "Objet:", avec ou sans accent/majuscule
        m = re.match(r"(?i)^objet\s*:\s*(.+)$", l)
        if m:
            return m.group(1).strip()
    # pas de champ Objet identifié : on renvoie le texte tel quel
    return llm_resume.strip()


def _set_lien_url(shape, texte: str, url: str):
    """ Écrit un texte cliquable pointant vers une URL externe. """
    tf = shape.text_frame
    p = tf.paragraphs[0]
    # réutilise le 1er run pour garder le style
    if p.runs:
        run = p.runs[0]
        run.text = texte
        for r in p.runs[1:]:
            r.text = ""
    else:
        run = p.add_run()
        run.text = texte
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    if url:
        run.hyperlink.address = url


def remplir_article(slide, art, trimestre=None):
    t = _shape_by_name(slide, "Titre 6")
    if t:
        _set_text(t, art.get("titre") or "Sans titre")
    # corps : uniquement l'Objet extrait du résumé structuré (concis, pas de débordement)
    corps = _shape_by_name(slide, "Rectangle 3")
    if corps:
        objet = _extraire_objet(art.get("llm_resume") or "")
        _set_text(corps, objet or art.get("resume") or "")
    # zones secondaires vidées pour éviter les chevauchements
    zt = _shape_by_name(slide, "ZoneTexte 9")
    if zt:
        _set_text(zt, "")
    zt2 = _shape_by_name(slide, "ZoneTexte 11")
    if zt2:
        meta = []
        if art.get("source_nom"):
            meta.append(art["source_nom"])
        if art.get("date_publication"):
            meta.append(art["date_publication"])
        _set_text(zt2, "  ·  ".join(meta))
    # source : URL cliquable
    src = _shape_by_name(slide, "Rectangle 16")
    if src:
        url = art.get("url") or ""
        _set_lien_url(src, url or "Source indisponible", url)
    if trimestre:
        _maj_trimestre(slide, trimestre)


def _maj_trimestre(slide, trimestre):
    for sh in slide.shapes:
        if sh.has_text_frame and "Trimestre" in sh.text_frame.text:
            _set_text(sh, trimestre)


def _lien_vers_slide(run, slide_cible, prs):
    """ Fait pointer un run de texte vers une autre slide (lien interne). """
    from pptx.oxml.ns import qn
    # partname de la slide cible, ex. /ppt/slides/slide4.xml
    rId = slide_cible.part.partname
    # créer la relation depuis la slide qui contient le run vers la slide cible
    source_part = run.part
    rel_id = source_part.relate_to(
        slide_cible.part,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
    )
    # attacher hlinkClick au run, avec action "jump to slide"
    rPr = run._r.get_or_add_rPr()
    # retirer un éventuel hlink existant
    for tag in ("a:hlinkClick",):
        ex = rPr.find(qn(tag))
        if ex is not None:
            rPr.remove(ex)
    hlink = rPr.makeelement(qn("a:hlinkClick"), {
        qn("r:id"): rel_id,
        "action": "ppaction://hlinksldjump",
    })
    # hlinkClick doit être en début de rPr
    rPr.insert(0, hlink)


def remplir_sommaire(slide, articles, slides_articles, prs):
    """ Remplit le sommaire ET crée les liens cliquables vers chaque article. """
    contenu = _shape_by_name(slide, "Espace réservé du contenu 4")
    if not contenu:
        return
    tf = contenu.text_frame
    modele = tf.paragraphs[0]
    style_run = modele.runs[0] if modele.runs else None

    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)

    if not articles:
        if style_run:
            style_run.text = "—"
        return

    # 1ère ligne : réutilise le paragraphe modèle
    if style_run:
        style_run.text = articles[0].get("titre", "")[:90]
        _lien_vers_slide(style_run, slides_articles[0], prs)

    # lignes suivantes
    for art, sl in zip(articles[1:], slides_articles[1:]):
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = art.get("titre", "")[:90]
        _lien_vers_slide(r, sl, prs)


def main(top=8, seuil=6, trimestre=None):
    parties = collecter(top, seuil)
    prs = Presentation(TEMPLATE)

    # mémorise les slides articles créées par partie (pour les liens du sommaire)
    slides_par_partie = {1: [], 2: [], 3: []}

    for p in (3, 2, 1):
        arts = parties[p]
        art_idx = ARTICLE[p]
        insert_pos = SOMMAIRE[p] + 1

        if not arts:
            _remove_slide(prs, art_idx)
            continue

        # 1er article dans la slide-type
        remplir_article(prs.slides[art_idx], arts[0], trimestre)
        slides_par_partie[p].append(prs.slides[art_idx])

        # duplication pour les suivants
        for i, art in enumerate(arts[1:], 1):
            new = _dup_slide(prs, art_idx)
            remplir_article(new, art, trimestre)
            _xml_move(prs, len(prs.slides) - 1, insert_pos + i)
            slides_par_partie[p].append(new)

        if trimestre:
            _maj_trimestre(prs.slides[SEP[p]], trimestre)
            _maj_trimestre(prs.slides[SOMMAIRE[p]], trimestre)

    # les sommaires + liens EN DERNIER, quand toutes les slides sont placées.
    # On retrouve chaque sommaire par son contenu (placeholder "Article 1..."),
    # car les index ont bougé avec les duplications.
    sommaires_restants = []
    for idx, s in enumerate(prs.slides):
        c = _shape_by_name(s, "Espace réservé du contenu 4")
        titre = _shape_by_name(s, "Titre 15")
        if c is not None and titre is not None and "Sommaire" in titre.text_frame.text:
            sommaires_restants.append(s)

    # les sommaires sont dans l'ordre des parties 1,2,3
    for p, som_slide in zip((1, 2, 3), sommaires_restants):
        if parties[p]:
            remplir_sommaire(som_slide, parties[p], slides_par_partie[p], prs)

    if trimestre:
        for sh in prs.slides[0].shapes:
            if sh.has_text_frame and "Trimestre" in sh.text_frame.text:
                _set_text(sh, trimestre)

    prs.save(SORTIE)
    for p in (1, 2, 3):
        print(f"Partie {p} : {len(parties[p])} articles")
    print(f"✓ {SORTIE} généré")

    # trimestre sur la couverture
    if trimestre:
        for sh in prs.slides[0].shapes:
            if sh.has_text_frame and "Trimestre" in sh.text_frame.text:
                _set_text(sh, trimestre)

    prs.save(SORTIE)
    for p in (1, 2, 3):
        print(f"Partie {p} : {len(parties[p])} articles")
    print(f"✓ {SORTIE} généré")


def _remove_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    sldIdLst.remove(ids[index])


if __name__ == "__main__":
    top, seuil, trimestre = 8, 6, None
    if "--top" in sys.argv:
        top = int(sys.argv[sys.argv.index("--top") + 1])
    if "--seuil" in sys.argv:
        seuil = int(sys.argv[sys.argv.index("--seuil") + 1])
    if "--trimestre" in sys.argv:
        trimestre = sys.argv[sys.argv.index("--trimestre") + 1]
    main(top, seuil, trimestre)
