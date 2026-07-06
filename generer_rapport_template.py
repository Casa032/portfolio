"""
generer_rapport_template.py — Remplit le template Cofidis avec les articles de la base.

Structure du template (11 slides) :
  1  couverture
  2  séparateur Partie 01 (FR)      3 sommaire      4 article
  5  séparateur Partie 02 (Intl)    6 sommaire      7 article
  8  séparateur Partie 03 (Juri)    9 sommaire     10 article
  11 contacts

Pour chaque partie : on garde le séparateur + le sommaire, et on duplique
la slide "article" autant de fois qu'il y a d'articles retenus.

Usage : python generer_rapport_template.py [--top N] [--seuil S] [--trimestre "2026 – 2e Trimestre"]
"""
import copy
import sys

from pptx import Presentation
from pptx.util import Pt

from database import get_connection
from classer_rapport import classer

TEMPLATE = "template.pptx"
SORTIE = "rapport_veille_cofidis.pptx"

# index (0-based) des slides du template
SEP = {1: 1, 2: 4, 3: 7}       # séparateur de chaque partie
SOMMAIRE = {1: 2, 2: 5, 3: 8}  # sommaire de chaque partie
ARTICLE = {1: 3, 2: 6, 3: 9}   # slide article-type de chaque partie


def meilleur_score(a):
    return a["llm_score"] if a.get("llm_score") is not None else (a.get("score_pertinence") or 0)


def collecter(top, seuil):
    conn = get_connection()
    try:
        rows = [dict(r) for r in conn.execute(
            "SELECT a.*, s.nom AS source_nom FROM articles a JOIN sources s ON s.id=a.source_id")]
    finally:
        conn.close()
    retenus = [a for a in rows if
               (a.get("llm_score") is not None and a["llm_score"] >= seuil)
               or (a.get("llm_score") is None and (a.get("score_pertinence") or 0) > 0)]
    parties = {1: [], 2: [], 3: []}
    for a in retenus:
        parties[classer(a)].append(a)
    for p in parties:
        parties[p].sort(key=meilleur_score, reverse=True)
        parties[p] = parties[p][:top]
    return parties


# ---- helpers python-pptx pour dupliquer/supprimer des slides ----

def _dup_slide(prs, index):
    """ Duplique la slide à l'index donné, l'ajoute en fin, renvoie la nouvelle. """
    source = prs.slides[index]
    blank = source.slide_layout
    new = prs.slides.add_slide(blank)
    # vider les placeholders ajoutés par le layout
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    # copier toutes les shapes de la source
    for sh in source.shapes:
        new.shapes._spTree.append(copy.deepcopy(sh._element))
    return new


def _xml_move(prs, from_pos, to_pos):
    """ Déplace la slide de from_pos vers to_pos dans l'ordre. """
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    el = ids[from_pos]
    sldIdLst.remove(el)
    sldIdLst.insert(to_pos, el)


def _set_text(shape, text):
    """ Remplace le texte d'une shape en gardant le style du 1er run. """
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r.text = ""
    else:
        p.add_run().text = text
    # supprimer les paragraphes suivants
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)


def _shape_by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def remplir_article(slide, art, trimestre=None):
    t = _shape_by_name(slide, "Titre 6")
    if t:
        _set_text(t, art.get("titre") or "Sans titre")
    # corps principal : résumé/raison dans Rectangle 3
    corps = _shape_by_name(slide, "Rectangle 3")
    if corps:
        _set_text(corps, art.get("resume") or art.get("llm_raison") or "")
    # zone secondaire : la raison LLM comme analyse, si dispo
    zt = _shape_by_name(slide, "ZoneTexte 9")
    if zt:
        raison = art.get("llm_raison") or ""
        _set_text(zt, f"Pertinence : {raison}" if raison else "")
    zt2 = _shape_by_name(slide, "ZoneTexte 11")
    if zt2:
        meta = []
        if art.get("source_nom"):
            meta.append(art["source_nom"])
        if art.get("date_publication"):
            meta.append(art["date_publication"])
        _set_text(zt2, "  ·  ".join(meta))
    src = _shape_by_name(slide, "Rectangle 16")
    if src:
        _set_text(src, f"Source : {art.get('url') or ''}")
    if trimestre:
        _maj_trimestre(slide, trimestre)


def _maj_trimestre(slide, trimestre):
    for sh in slide.shapes:
        if sh.has_text_frame and "Trimestre" in sh.text_frame.text:
            _set_text(sh, trimestre)


def _lien_vers_slide(run, slide_cible, prs):
    """ Fait pointer un run de texte vers une autre slide (lien interne). """
    from pptx.oxml.ns import qn
    # partname de la slide cible, ex. /ppt/slides/slide4.xml
    rId = slide_cible.part.partname
    # créer la relation depuis la slide qui contient le run vers la slide cible
    source_part = run.part
    rel_id = source_part.relate_to(
        slide_cible.part,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
    )
    # attacher hlinkClick au run, avec action "jump to slide"
    rPr = run._r.get_or_add_rPr()
    # retirer un éventuel hlink existant
    for tag in ("a:hlinkClick",):
        ex = rPr.find(qn(tag))
        if ex is not None:
            rPr.remove(ex)
    hlink = rPr.makeelement(qn("a:hlinkClick"), {
        qn("r:id"): rel_id,
        "action": "ppaction://hlinksldjump",
    })
    # hlinkClick doit être en début de rPr
    rPr.insert(0, hlink)


def remplir_sommaire(slide, articles, slides_articles, prs):
    """ Remplit le sommaire ET crée les liens cliquables vers chaque article. """
    contenu = _shape_by_name(slide, "Espace réservé du contenu 4")
    if not contenu:
        return
    tf = contenu.text_frame
    modele = tf.paragraphs[0]
    style_run = modele.runs[0] if modele.runs else None

    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)

    if not articles:
        if style_run:
            style_run.text = "—"
        return

    # 1ère ligne : réutilise le paragraphe modèle
    if style_run:
        style_run.text = articles[0].get("titre", "")[:90]
        _lien_vers_slide(style_run, slides_articles[0], prs)

    # lignes suivantes
    for art, sl in zip(articles[1:], slides_articles[1:]):
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = art.get("titre", "")[:90]
        _lien_vers_slide(r, sl, prs)


def main(top=8, seuil=6, trimestre=None):
    parties = collecter(top, seuil)
    prs = Presentation(TEMPLATE)

    # mémorise les slides articles créées par partie (pour les liens du sommaire)
    slides_par_partie = {1: [], 2: [], 3: []}

    for p in (3, 2, 1):
        arts = parties[p]
        art_idx = ARTICLE[p]
        insert_pos = SOMMAIRE[p] + 1

        if not arts:
            _remove_slide(prs, art_idx)
            continue

        # 1er article dans la slide-type
        remplir_article(prs.slides[art_idx], arts[0], trimestre)
        slides_par_partie[p].append(prs.slides[art_idx])

        # duplication pour les suivants
        for i, art in enumerate(arts[1:], 1):
            new = _dup_slide(prs, art_idx)
            remplir_article(new, art, trimestre)
            _xml_move(prs, len(prs.slides) - 1, insert_pos + i)
            slides_par_partie[p].append(new)

        if trimestre:
            _maj_trimestre(prs.slides[SEP[p]], trimestre)
            _maj_trimestre(prs.slides[SOMMAIRE[p]], trimestre)

    # les sommaires + liens EN DERNIER, quand toutes les slides sont placées.
    # On retrouve chaque sommaire par son contenu (placeholder "Article 1..."),
    # car les index ont bougé avec les duplications.
    sommaires_restants = []
    for idx, s in enumerate(prs.slides):
        c = _shape_by_name(s, "Espace réservé du contenu 4")
        titre = _shape_by_name(s, "Titre 15")
        if c is not None and titre is not None and "Sommaire" in titre.text_frame.text:
            sommaires_restants.append(s)

    # les sommaires sont dans l'ordre des parties 1,2,3
    for p, som_slide in zip((1, 2, 3), sommaires_restants):
        if parties[p]:
            remplir_sommaire(som_slide, parties[p], slides_par_partie[p], prs)

    if trimestre:
        for sh in prs.slides[0].shapes:
            if sh.has_text_frame and "Trimestre" in sh.text_frame.text:
                _set_text(sh, trimestre)

    prs.save(SORTIE)
    for p in (1, 2, 3):
        print(f"Partie {p} : {len(parties[p])} articles")
    print(f"✓ {SORTIE} généré")

    # trimestre sur la couverture
    if trimestre:
        for sh in prs.slides[0].shapes:
            if sh.has_text_frame and "Trimestre" in sh.text_frame.text:
                _set_text(sh, trimestre)

    prs.save(SORTIE)
    for p in (1, 2, 3):
        print(f"Partie {p} : {len(parties[p])} articles")
    print(f"✓ {SORTIE} généré")


def _remove_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    sldIdLst.remove(ids[index])


if __name__ == "__main__":
    top, seuil, trimestre = 8, 6, None
    if "--top" in sys.argv:
        top = int(sys.argv[sys.argv.index("--top") + 1])
    if "--seuil" in sys.argv:
        seuil = int(sys.argv[sys.argv.index("--seuil") + 1])
    if "--trimestre" in sys.argv:
        trimestre = sys.argv[sys.argv.index("--trimestre") + 1]
    main(top, seuil, trimestre)
